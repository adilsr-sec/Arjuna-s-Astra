from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math

# ── Palette ──────────────────────────────────────────────────────────────────
BG    = RGBColor(0x0A, 0x12, 0x0A)
GRN   = RGBColor(0x4A, 0xDE, 0x80)
AMB   = RGBColor(0xFF, 0xB7, 0x03)
WHT   = RGBColor(0xFF, 0xFF, 0xFF)
GRY   = RGBColor(0x88, 0x99, 0x88)
RED   = RGBColor(0xFF, 0x44, 0x44)
BLU   = RGBColor(0x38, 0xBD, 0xF8)
DRK   = RGBColor(0x06, 0x0C, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    return s

def rect(s, l, t, w, h, fill=GRN, line_color=None, line_w=Pt(1)):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def hrule(s, l, t, w, color=GRN, thick=0.04):
    rect(s, l, t, w, thick, color)

def txt(s, l, t, w, h, text, size=13, bold=False, color=WHT,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Courier New"
    return tb

def bullets(s, l, t, w, h, title, items, tc=GRN, ic=WHT, ts=14, is_=12):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = title; r.font.size = Pt(ts); r.font.bold = True
    r.font.color.rgb = tc; r.font.name = "Courier New"
    for item in items:
        p2 = tf.add_paragraph(); p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = f"  ▸  {item}"
        r2.font.size = Pt(is_); r2.font.color.rgb = ic
        r2.font.name = "Courier New"

def header(s, title, subtitle="", slide_num=""):
    hrule(s, 0, 0, 13.33, AMB, 0.07)
    txt(s, 0.4, 0.12, 10, 0.55, title, 22, True, GRN)
    if subtitle:
        txt(s, 0.4, 0.68, 12, 0.4, subtitle, 12, False, GRY)
    if slide_num:
        txt(s, 12.3, 0.15, 0.9, 0.4, slide_num, 10, False, AMB, PP_ALIGN.RIGHT)
    hrule(s, 0, 7.3, 13.33, GRN, 0.06)

def flow_boxes(s, labels, colors, t=3.5, box_w=1.6, box_h=0.65):
    n = len(labels)
    total = n * box_w + (n-1) * 0.45
    start_l = (13.33 - total) / 2
    for i, (lbl, col) in enumerate(zip(labels, colors)):
        l = start_l + i * (box_w + 0.45)
        sh = rect(s, l, t, box_w, box_h, col)
        txt(s, l, t, box_w, box_h, lbl, 10, True, DRK, PP_ALIGN.CENTER)
        if i < n - 1:
            txt(s, l + box_w + 0.02, t + 0.18, 0.41, 0.3, ">>>", 11, True, AMB)

def table_slide(s, headers, rows, t=1.3, col_widths=None):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [12.5 / n_cols] * n_cols
    # header row
    l = 0.4
    for i, h in enumerate(headers):
        rect(s, l, t, col_widths[i], 0.42, AMB)
        txt(s, l+0.05, t+0.05, col_widths[i]-0.1, 0.35, h, 11, True, DRK)
        l += col_widths[i]
    # data rows
    for ri, row in enumerate(rows):
        l = 0.4
        bg_c = RGBColor(0x12, 0x20, 0x12) if ri % 2 == 0 else RGBColor(0x0E, 0x18, 0x0E)
        for ci, cell in enumerate(row):
            rect(s, l, t + 0.42 + ri * 0.52, col_widths[ci], 0.52, bg_c, GRN)
            txt(s, l+0.05, t + 0.42 + ri * 0.52 + 0.08,
                col_widths[ci]-0.1, 0.4, cell, 10, False, WHT)
            l += col_widths[ci]

def arch_box(s, l, t, w, h, label, sublabel="", fill=GRN, tc=DRK):
    rect(s, l, t, w, h, fill, AMB)
    txt(s, l+0.05, t+0.08, w-0.1, 0.3, label, 10, True, tc, PP_ALIGN.CENTER)
    if sublabel:
        txt(s, l+0.05, t+0.35, w-0.1, 0.3, sublabel, 8, False, DRK, PP_ALIGN.CENTER)
