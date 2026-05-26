from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color Palette - Dark Military Theme
BG_COLOR   = RGBColor(0x0D, 0x14, 0x0D)  # Very dark green-black
ACCENT1    = RGBColor(0x4A, 0xDE, 0x80)  # Tactical green
ACCENT2    = RGBColor(0xFF, 0xB7, 0x03)  # Amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xA0, 0xB0, 0xA0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Completely blank layout

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg_color=None, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Courier New"
    return txBox

def line(slide, l, t, w, color=ACCENT1, thickness=36000):
    """Draw a horizontal rule using a rectangle."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Emu(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def bullet_block(slide, l, t, w, h, title, items, title_color=ACCENT1,
                 item_color=WHITE, item_size=13):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    # Title paragraph
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Courier New"
    # Bullet items
    for item in items:
        para = tf.add_paragraph()
        para.space_before = Pt(3)
        run2 = para.add_run()
        run2.text = f"  > {item}"
        run2.font.size = Pt(item_size)
        run2.font.color.rgb = item_color
        run2.font.name = "Courier New"

# ─── Slide 1: Title ───────────────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT2, 72000)
box(s, 0.5, 0.6, 12.3, 1.2, "ARJUNA'S ARROW", 54, True, ACCENT1, PP_ALIGN.CENTER)
box(s, 0.5, 1.7, 12.3, 0.6, "PROJECT STEG — COVERT INTELLIGENCE TRANSMISSION SYSTEM",
    18, False, ACCENT2, PP_ALIGN.CENTER)
line(s, 1, 2.5, 11.33, ACCENT1, 36000)
box(s, 0.5, 2.8, 12.3, 0.5, "Military-Grade Steganographic Communication Platform",
    15, False, GREY, PP_ALIGN.CENTER)
box(s, 0.5, 3.3, 12.3, 0.5,
    "AES-256 Encryption  |  RSA Key Exchange  |  LSB Steganography  |  ARSS Audio",
    14, False, GREY, PP_ALIGN.CENTER)
line(s, 0, 7.2, 13.33, ACCENT2, 54000)
box(s, 0.5, 6.8, 12.3, 0.4, "[ CLASSIFIED — FOR INTERNAL REVIEW ONLY ]",
    11, False, ACCENT2, PP_ALIGN.CENTER)

# ─── Slide 2: Project Overview ────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 6, 0.5, "01 // PROJECT OVERVIEW", 20, True, ACCENT1)
box(s, 0.4, 0.85, 12.3, 1.6,
    "ARJUNA'S ARROW is a military-grade covert communication platform. "
    "It allows field operatives to encode classified audio intelligence "
    "into ordinary-looking image files using steganography, then transmit "
    "them securely across a network — undetectable to any outside observer.",
    14, False, WHITE)
line(s, 0.4, 2.55, 12.5, ACCENT2, 36000)
bullet_block(s, 0.4, 2.75, 5.7, 4.2, "PROBLEM STATEMENT", [
    "Standard networks are monitored & intercepted",
    "Encrypted files raise immediate suspicion",
    "Voice intelligence is high-bandwidth & traceable",
    "No secure, deniable audio comm channel existed",
])
bullet_block(s, 6.8, 2.75, 6, 4.2, "OUR SOLUTION", [
    "Audio → Spectrogram (ARSS compression)",
    "Spectrogram → AES-256 + RSA encryption",
    "Encrypted data → Hidden in image pixels (LSB)",
    "Image looks normal — zero visual difference",
    "Recipient reverses process to get audio back",
], title_color=ACCENT2)

# ─── Slide 3: Tech Stack ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 8, 0.5, "02 // SYSTEM ARCHITECTURE & TECH STACK", 20, True, ACCENT1)

bullet_block(s, 0.4, 0.9, 5.8, 3.2, "FRONTEND", [
    "Framework: Next.js (React 18)",
    "Styling: Tailwind CSS",
    "Auth State: JWT + localStorage session",
    "Audio: Web Audio API (live mic recording)",
    "Design: Dark Olive & Amber Military Theme",
])
bullet_block(s, 6.8, 0.9, 6, 3.2, "BACKEND", [
    "Framework: FastAPI (Python)",
    "Database: SQLite + SQLAlchemy ORM",
    "Auth: PyJWT + Passlib (PBKDF2-SHA256)",
    "Crypto: Python cryptography (AES-GCM + RSA)",
    "Audio: librosa + soundfile + numpy",
    "Stego: OpenCV (cv2) + numpy",
])
line(s, 0.4, 4.2, 12.5, ACCENT2, 36000)
box(s, 0.4, 4.4, 12.3, 0.4, "DEPLOYMENT: Single-machine, two-process. Backend on port 8000 (Uvicorn). Frontend on port 3000 (Next.js Dev Server). Launched via START_ARJUNAS_ARROW.bat", 12, False, GREY)

# ─── Slide 4: Core Pipeline ───────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 10, 0.5, "03 // TRANSMISSION PIPELINE (4 PHASES)", 20, True, ACCENT1)

phases = [
    ("PHASE 1", "ARSS AUDIO CONVERSION",
     "Audio (.wav/.mp3 or live mic) is loaded by librosa. A Mel-Spectrogram "
     "(96 mels, 1024-point FFT) converts waveform to a compact frequency matrix. "
     "OpenCV generates a visual heatmap for UI preview."),
    ("PHASE 2", "HYBRID CRYPTOGRAPHY",
     "Spectrogram data is encrypted with AES-256-GCM (symmetric). The AES key "
     "is then wrapped with the recipient's RSA Public Key (asymmetric). Optional "
     "PBKDF2 password adds a second authentication layer."),
    ("PHASE 3", "LSB STEGANOGRAPHY",
     "steg.py uses a seed-initialized PRNG (random.Random(seed)) to shuffle "
     "pixel positions in the cover image. It overwrites the Least Significant "
     "Bit of each channel with payload bits. PSNR > 40dB — visually identical."),
    ("PHASE 4", "SECURE TRANSMISSION",
     "The stego image is stored on the server and linked to the recipient's "
     "inbox as an 'encoded_transmission'. To any network observer it appears "
     "as routine image traffic. Files are scheduled for deletion post-download."),
]

for i, (tag, title, desc) in enumerate(phases):
    col = i * 3.3 + 0.4
    box(s, col, 0.9, 3.0, 0.35, tag, 10, True, ACCENT2)
    box(s, col, 1.2, 3.0, 0.4, title, 13, True, ACCENT1)
    line(s, col, 1.65, 3.0, ACCENT1, 27000)
    box(s, col, 1.8, 3.0, 4.0, desc, 11, False, WHITE)

line(s, 0.4, 6.3, 12.5, ACCENT2, 27000)
box(s, 0.4, 6.45, 12.3, 0.4,
    "Audio  >>  ARSS Spectrogram  >>  AES-256/RSA Encrypt  >>  LSB Embed in Image  >>  Network Send",
    13, True, ACCENT2, PP_ALIGN.CENTER)

# ─── Slide 5: RBAC & Modules ──────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 10, 0.5, "04 // ROLE-BASED ACCESS CONTROL (RBAC)", 20, True, ACCENT1)

bullet_block(s, 0.4, 0.9, 3.9, 5.8, "ADMIN (Root Clearance)", [
    "Provision Departments & Soldiers",
    "Hierarchy Management (Edit/Delete any user)",
    "Terminal Re-keying (Password Override)",
    "System Audit Logs (all events)",
    "Global Broadcast (encode to anyone)",
    "Network Intercepts (read ALL traffic)",
], title_color=RGBColor(0xFF,0x44,0x44))

bullet_block(s, 4.7, 0.9, 3.9, 5.8, "DEPARTMENT (Sector Command)", [
    "Provision Soldiers (own dept only)",
    "Operatives Roster (view/edit/delete unit)",
    "Cannot access other departments",
    "Send & receive encoded intel",
    "Decode intercepted stego images",
    "Decode signals (spectrogram viewer)",
], title_color=ACCENT2)

bullet_block(s, 9.0, 0.9, 3.9, 5.8, "SOLDIER (Field Operative)", [
    "Live mic audio recording",
    "Encode & transmit stego images",
    "Receive & decode stego messages",
    "Griffin-Lim audio reconstruction",
    "Tactical HUD dashboard",
    "Emergency Burn Protocol button",
], title_color=ACCENT1)

# ─── Slide 6: Security Features ───────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 10, 0.5, "05 // SECURITY FEATURES & DESIGN DECISIONS", 20, True, ACCENT1)

features = [
    ("ZERO-KNOWLEDGE PASSWORDS",
     "All passwords hashed with PBKDF2-SHA256 via Passlib. Cannot be reversed or viewed — not even by the database admin."),
    ("SEED-LOCKED STEGANOGRAPHY",
     "Even if the stego image is intercepted, the payload cannot be extracted without the exact numeric seed that randomizes pixel order."),
    ("CASCADING DATA PURGE",
     "Deleting an operative triggers automatic deletion of all their messages and telemetry logs, eliminating database traces."),
    ("SELF-DESTRUCTING MEDIA",
     "Decoded audio files & spectrograms are served once then scheduled for server-side deletion via a background cleanup task."),
    ("TERMINAL RE-KEYING",
     "Since passwords cannot be decrypted, Admins override compromised credentials by re-keying — the new hash replaces the old one securely."),
    ("RSA KEY-PER-USER",
     "Each user has their own RSA key pair. Messages can only be decrypted by the intended recipient using their private key."),
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col_idx = i % 2
    l = 0.4 + col_idx * 6.4
    t = 1.0 + row * 2.0
    box(s, l, t, 6.0, 0.35, title, 12, True, ACCENT2)
    box(s, l, t + 0.38, 6.0, 1.4, desc, 12, False, WHITE)

# ─── Slide 7: UI Modules ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 10, 0.5, "06 // DASHBOARD MODULES & UI FEATURES", 20, True, ACCENT1)

bullet_block(s, 0.4, 0.9, 5.8, 6.2, "ADMIN DASHBOARD", [
    "Global Command & Control header",
    "Provision Department / Soldier panels",
    "Hierarchy Management (tree view)",
    "Edit Operative modal with Re-keying",
    "Decode Intercepted Signal panel",
    "Global Broadcast (encode + live audio)",
    "Global Comms & Intercepts (full wiretap)",
    "System Audit Logs (auto-refreshed)",
    "Spectrogram viewer modal",
])

bullet_block(s, 6.8, 0.9, 6.0, 3.0, "DEPARTMENT DASHBOARD", [
    "C2 Command Center header",
    "Operatives Roster (unit management)",
    "Encode & send intel panel",
    "Decode Intercepted Signal panel",
    "Department inbox viewer",
])

bullet_block(s, 6.8, 4.1, 6.0, 3.2, "SOLDIER TERMINAL", [
    "Mission Objectives widget (tactical)",
    "System Integrity monitor (live values)",
    "Terminal Telemetry (DEFCON, GPS, time)",
    "Security Radar sweep animation",
    "Transmission Pipeline flow display",
    "Live audio recorder + file upload",
    "Emergency Burn Protocol button",
])

# ─── Slide 8: Database Schema ─────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT1, 54000)
box(s, 0.4, 0.15, 10, 0.5, "07 // DATABASE SCHEMA", 20, True, ACCENT1)

bullet_block(s, 0.4, 0.9, 3.8, 5.5, "TABLE: users", [
    "id (PK)",
    "username (unique, lowercase)",
    "hashed_password (PBKDF2)",
    "role: Admin / Department / Soldier",
    "department (sector name)",
    "public_key (RSA PEM)",
    "locked (boolean — burn protocol)",
], title_color=ACCENT2)

bullet_block(s, 4.6, 0.9, 3.8, 5.5, "TABLE: messages", [
    "id (PK)",
    "sender_id (FK → users)",
    "recipient_id (FK → users)",
    "body (text payload / status)",
    "type: text / encoded_transmission",
    "stego_image (server file path)",
    "ts (timestamp)",
], title_color=ACCENT2)

bullet_block(s, 9.0, 0.9, 3.8, 5.5, "TABLE: telemetry_logs", [
    "id (PK)",
    "event (LOGIN, ENCODE, DELETE...)",
    "actor (username)",
    "status (SUCCESS / FAIL)",
    "details (description)",
    "ts (timestamp)",
], title_color=ACCENT2)

line(s, 0.4, 6.5, 12.5, ACCENT2, 27000)
box(s, 0.4, 6.6, 12.3, 0.4, "SQL: SELECT DISTINCT department FROM users WHERE role = 'Department';", 12, True, ACCENT1, PP_ALIGN.CENTER)

# ─── Slide 9: Summary ─────────────────────────────────────────────────────────
s = add_slide(); bg(s)
line(s, 0, 0, 13.33, ACCENT2, 72000)
box(s, 0.5, 0.6, 12.3, 0.6, "08 // SUMMARY", 22, True, ACCENT1, PP_ALIGN.CENTER)
line(s, 1, 1.4, 11.33, ACCENT1, 36000)

summary_points = [
    "Full-stack covert communication platform (Next.js + FastAPI + SQLite)",
    "4-phase Transmission Pipeline: ARSS Audio > AES/RSA Encrypt > LSB Stego > Network Send",
    "3-tier RBAC: Admin (global) > Department (sector) > Soldier (operative)",
    "Admin can intercept ALL network traffic and override any credentials",
    "Zero-Knowledge password storage — mathematically unrecoverable",
    "PRNG Seed-locked steganography — image appears visually identical to original",
    "Cascading data purge on operative deletion — no orphaned traces",
    "Self-destructing decoded media — files deleted after single download",
    "Live browser microphone recording — no external software needed",
]
for i, point in enumerate(summary_points):
    box(s, 0.8, 1.6 + i * 0.52, 11.8, 0.5, f"[{i+1:02d}]  {point}", 13, False, WHITE if i % 2 == 0 else GREY)

line(s, 0, 7.2, 13.33, ACCENT2, 54000)
box(s, 0.5, 6.8, 12.3, 0.4, "[ END OF CLASSIFIED BRIEFING — ARJUNA'S ARROW ]", 13, True, ACCENT2, PP_ALIGN.CENTER)

# Save
out_path = r"C:\Users\adhil\OneDrive\Desktop\Project\STEG\Arjunas_Arrow_Presentation.pptx"
prs.save(out_path)
print(f"PPT saved to: {out_path}")
